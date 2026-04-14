"""
PowerPoint-ready chart export utilities.

Standard PowerPoint widescreen slide: 13.33 in × 7.5 in.
For a full-bleed chart on a 16:9 slide at 96 dpi: 1280 × 720 px.
We target 10 in × 5.6 in (content area) so text/axes stay comfortable.
"""

from __future__ import annotations

import io
import os
from typing import Optional

import plotly.graph_objects as go

# ---------------------------------------------------------------------------
# Dimensions for PowerPoint export (inches × dpi → pixels)
# ---------------------------------------------------------------------------
PPT_WIDTH_IN = 10.0
PPT_HEIGHT_IN = 5.6
PPT_DPI = 144          # high-res for crisp export
PPT_WIDTH_PX = int(PPT_WIDTH_IN * PPT_DPI)    # 1440
PPT_HEIGHT_PX = int(PPT_HEIGHT_IN * PPT_DPI)  # 806

# Plotly defaults for screen display
SCREEN_WIDTH = 900
SCREEN_HEIGHT = 480


def ppt_layout(title: str = "", font_size: int = 14) -> dict:
    """Return a Plotly layout dict configured for PowerPoint export."""
    return dict(
        title=dict(text=title, font=dict(size=font_size + 2)),
        font=dict(family="Calibri, Arial, sans-serif", size=font_size),
        plot_bgcolor="white",
        paper_bgcolor="white",
        margin=dict(l=60, r=40, t=60, b=60),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
    )


def fig_to_png_bytes(fig: go.Figure, ppt: bool = False) -> bytes:
    """
    Render a Plotly figure to PNG bytes.
    If ppt=True, use PowerPoint dimensions; otherwise screen dimensions.
    """
    w = PPT_WIDTH_PX if ppt else SCREEN_WIDTH
    h = PPT_HEIGHT_PX if ppt else SCREEN_HEIGHT
    try:
        return fig.to_image(format="png", width=w, height=h, scale=1)
    except Exception:
        # kaleido not installed or failed; return empty bytes
        return b""


def build_pptx(figures: list[dict]) -> bytes:
    """
    Build a PPTX file from a list of figure dicts.

    Each dict should have:
      - fig: plotly Figure
      - title: slide title string
      - notes: speaker notes string (optional)

    Returns raw PPTX bytes.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank slide

    for item in figures:
        fig: go.Figure = item["fig"]
        slide_title: str = item.get("title", "")
        notes_text: str = item.get("notes", "")

        # Update figure for PPT dimensions
        fig_ppt = go.Figure(fig)
        fig_ppt.update_layout(
            width=PPT_WIDTH_PX,
            height=PPT_HEIGHT_PX,
            **ppt_layout(slide_title),
        )

        img_bytes = fig_to_png_bytes(fig_ppt, ppt=True)

        slide = prs.slides.add_slide(blank_layout)

        # Title text box
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.55)
        )
        tf = title_box.text_frame
        tf.text = slide_title
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True

        # Image (charts)
        if img_bytes:
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(
                img_stream,
                Inches(0.4),
                Inches(0.85),
                Inches(PPT_WIDTH_IN),
                Inches(PPT_HEIGHT_IN),
            )

        # Notes
        if notes_text:
            notes_slide = slide.notes_slide
            tf_notes = notes_slide.notes_text_frame
            tf_notes.text = notes_text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def fig_to_svg_bytes(fig: go.Figure) -> bytes:
    """Return SVG bytes for a figure (vector, resolution-independent)."""
    try:
        return fig.to_image(format="svg", width=PPT_WIDTH_PX, height=PPT_HEIGHT_PX)
    except Exception:
        return b""
